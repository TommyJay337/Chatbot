from openai import OpenAI
import os
from dotenv import load_dotenv
from pptx import Presentation
import PyPDF2
import json
from cassandra.cluster import Cluster
from cassandra.auth import PlainTextAuthProvider
from langchain.memory import CassandraChatMessageHistory, ConversationBufferMemory

load_dotenv()

# Open the JSON file for reading
with open("vector_database-token.json") as f:
    secrets = json.load(f)

auth_provider = PlainTextAuthProvider(username=secrets["clientId"], password=secrets["secret"])
cloud_config = {'secure_connect_bundle': 'secure-connect-vector-database.zip'}
cluster = Cluster(cloud=cloud_config, auth_provider=auth_provider)
session = cluster.connect()

ASTRA_DB_KEYSPACE = "search"
message_history = CassandraChatMessageHistory(
    session_id="session_identifier",  # This could be dynamically set per user or conversation
    session=session,
    keyspace=ASTRA_DB_KEYSPACE,
    ttl_seconds=3600  # Adjust as needed
)

cass_buff_memory = ConversationBufferMemory(
    memory_key="chat_history",
    chat_memory=message_history
)

# Optionally clear history if needed at the start of a new conversation
# message_history.clear()

def extract_text_from_multiple_pdfs(file_list):
    all_text_content = []
    for pdf_file in file_list:
        with open(pdf_file, "rb") as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    all_text_content.append(page_text.replace('\n', ' '))
    return ' '.join(all_text_content)

def process_pdfs_in_folder(folder_path):
    pdf_files = [os.path.join(folder_path, f) for f in os.listdir(folder_path) if f.lower().endswith('.pdf')]
    if len(pdf_files) > 0:
        extracted_text = extract_text_from_multiple_pdfs(pdf_files)
        return extracted_text
    else:
        return "Folder does not contain any PDF files."

def extract_text_from_multiple_pptx(file_list):
    all_text_content = []
    for pptx_file in file_list:
        prs = Presentation(pptx_file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_text = ' '.join(paragraph.text for paragraph in shape.text_frame.paragraphs)
                    all_text_content.append(slide_text)
    return ' '.join(all_text_content)

neuro_files = ["Lectures/neuro/ALS part 1.pptx",
                "Lectures/neuro/Alzheimer's 2024.pptx",
                "Lectures/neuro/Cognitive perceptual and behavioural considerations 2024.pptx",
                "Lectures/neuro/gbs.pptx",
                "Lectures/neuro/PPS.pptx"
                ]

sports_files = ["Lectures/sports/Concussion Lecture_2024.pptx",
                "Lectures/sports/HOPS + treatment.pptx",
                "Lectures/sports/Injury Prevention - School PT.pptx",
                "Lectures/sports/Intro to Sport Medicine lecture.pptm",
                #"Lectures/sports/Nutrition in Sport_2024.pdf",
                "Lectures/sports/Practical Application Primary_Secondary.pptx",
                #"Lectures/sports/Return to Performance, Post-Op Protocols.pdf",
                "Lectures/sports/Return to Sport 2024.pptx",
                "Lectures/sports/Sideline coverage_EAP, primary and secondary.pptx",
                "Lectures/sports/Sport Psychology_2024.pptx",
                "Lectures/sports/Sports Massage Therapy PT_s.pptx"
                ]

def choose_class_content(user_input):
    if user_input == 'neuro':
        return extract_text_from_multiple_pptx(neuro_files)
    elif user_input == 'sports':
        text_from_pptx = extract_text_from_multiple_pptx(sports_files)
        text_from_pdfs = process_pdfs_in_folder("Lectures/sports")
        return text_from_pptx + ' ' + text_from_pdfs
    else:
        return "Invalid input. Please choose 'neuro' or 'sports'."

user_input = input("Which class would you like to study today? ('neuro', 'sports'): ")

# Extract text from all specified files
context_text = choose_class_content(user_input)
#print(context_text)

api_key = os.getenv("OPENAI_API_KEY")

'''
if api_key:
    print("API key is loaded.")
else:
    print("API key not found. Please check your environment variables.")
    exit(1)
'''

client = OpenAI(api_key=api_key)
print("The OpenAI client has been successfully initialized.")

print("Chatbot: Hello! How would you like to study today?")

# Assuming we maintain our history list
conversation_history = []

conversation_history = []

def chat_with_gpt(prompt, context_text, instruction):
    global conversation_history

    # Append user's prompt to history
    conversation_history.append({"role": "user", "content": prompt})

    # Prepare the chat messages including context and instructions as part of the system messages
    messages = [
        {"role": "system", "content": context_text},
        {"role": "system", "content": instruction}
    ] + conversation_history

    # Use the new interface to create chat completions
    response = client.chat.completions.create(
        model="gpt-4-turbo",
        messages=messages,
        temperature=1,
        max_tokens=256,
        top_p=1,
        frequency_penalty=0,
        presence_penalty=0
    )

    # Extract the bot's response and append it to the history
    bot_response = response.choices[0].message.content.strip()
    conversation_history.append({"role": "assistant", "content": bot_response})

    # Return the content of the assistant's message
    return bot_response


if __name__ == "__main__":
    instruction = f"You are a helpful tutor for a physiotherapy student studying a neuro physiotherapy class and a sports physiotherapy class. You will explain concepts, clarify doubts, and assist in the understanding of topics. You will provide detailed explanations, generate flashcard-style quizzes with explanations for answers, and help with note-taking to compare and contrast the various conditions covered in the presentations. In addition to multiple choice questions, the tutor will also include case studies and short answer questions in quizzes to help students apply their knowledge in more practical scenarios. The tutor will ask for clarifications when needed and encourages the user to teach back the learned material to reinforce understanding."
    while True:
        user_input = input("You: ")
        if user_input.lower() in ['exit', 'quit', 'q', 'bye']:
            break
        response = chat_with_gpt(user_input, context_text, instruction)
        print("ChatBot: ", response)

