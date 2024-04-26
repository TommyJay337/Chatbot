import openai
import os
from dotenv import load_dotenv
from pptx import Presentation

load_dotenv()

def extract_text_from_pptx(pptx_file):
    text_content = []
    prs = Presentation(pptx_file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_content.append(' '.join(paragraph.text for paragraph in shape.text_frame.paragraphs))
    return ' '.join(text_content)

# Specify your PowerPoint file path
pptx_file = r'Lectures\gbs.pptx'
context_text = extract_text_from_pptx(pptx_file)

api_key = os.getenv("OPENAI_API_KEY")
if api_key:
    print("API key is loaded.")
else:
    print("API key not found. Please check your environment variables.")
    exit(1)

client = openai.OpenAI(api_key=api_key)
print("The OpenAI client has been successfully initialized.")

def chat_with_gpt(prompt, context_text, instruction):
    response = client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[
                {"role": "system", "content": context_text},  # Providing context
                {"role": "system", "content": instruction},  # Instruction prompt
                {"role": "user", "content": prompt}
                ],
                temperature=1,
                max_tokens=256,
                top_p=1,
                frequency_penalty=0,
                presence_penalty=0
            )
    return response.choices[0].message.content.strip()

if __name__ == "__main__":
    instruction = "You are a helpful tutor for a physiotherapy student who is studying Guillaine Barre Syndrome (GBS)."
    while True:
        user_input = input("You: ")
        if user_input.lower() in ['exit', 'quit', 'q', 'bye']:
            break
        response = chat_with_gpt(user_input, context_text, instruction)
        print("ChatBot: ", response)